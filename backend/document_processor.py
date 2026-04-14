from logger_config import get_logger

logger = get_logger(__name__)

"""
Document processing module for Thesis
Handles text extraction, chunking, and embedding generation
"""
import os
from io import BytesIO
from typing import Dict, List

import voyageai
from dotenv import load_dotenv

from api.utils.error_handler import (
    DocumentProcessingError,
    EmbeddingGenerationError,
    wrap_external_service_error,
)
from api.utils.retry import retry_with_backoff
import pb_client as pb
from cache import cache_search_results, get_cached_search_results
from config.constants import DATABASE, EMBEDDING, SEARCH, TEXT_CHUNKING
from repositories import conversations as conversations_repo
from repositories import documents as documents_repo

try:
    from openpyxl import load_workbook

    XLSX_SUPPORT = True
except ImportError:
    XLSX_SUPPORT = False
    logger.warning("openpyxl not installed - XLSX support disabled")

try:
    from pptx import Presentation

    PPTX_SUPPORT = True
except ImportError:
    PPTX_SUPPORT = False
    logger.warning("python-pptx not installed - PPTX support disabled")

# OCR support for image-based PDFs
try:
    import pytesseract
    from pdf2image import convert_from_bytes

    OCR_SUPPORT = True
except ImportError:
    OCR_SUPPORT = False
    logger.warning("pytesseract/pdf2image not installed - OCR support disabled")


def _extract_pdf_with_ocr(file_data: bytes, filename: str) -> str:
    """Extract text from a PDF using OCR (for image-based PDFs).

    Args:
        file_data: Raw PDF file bytes
        filename: Name of the file (for logging)

    Returns:
        Extracted text from all pages
    """
    if not OCR_SUPPORT:
        raise ValueError("OCR not available - install pytesseract and pdf2image")

    try:
        # Convert PDF pages to images
        images = convert_from_bytes(file_data, dpi=200)
        logger.info(f"   Converting {len(images)} pages to images for OCR...")

        text_parts = []
        for page_num, image in enumerate(images, 1):
            # Run OCR on each page
            text = pytesseract.image_to_string(image)
            if text.strip():
                text_parts.append(f"=== Page {page_num} ===\n{text}")

        full_text = "\n\n".join(text_parts)

        if not full_text.strip():
            raise ValueError("OCR could not extract any text from the PDF")

        logger.info(f"   OCR extracted {len(full_text)} characters from {filename}")
        return full_text

    except Exception as e:
        logger.error(f"❌ OCR failed for {filename}: {str(e)}")
        raise ValueError(f"OCR failed: {str(e)}") from None


load_dotenv()

# Initialize clients
VOYAGE_API_KEY = os.environ.get("VOYAGE_API_KEY")

vo = voyageai.Client(api_key=VOYAGE_API_KEY)


def extract_text_from_file(file_data: bytes, filename: str) -> str:
    """Extract text from different file types.

    Args:
        file_data: Raw file bytes
        filename: Name of the file (used to determine type)

    Returns:
        Extracted text content
    """
    file_ext = filename.lower().split(".")[-1]

    if file_ext == "docx":
        # Extract text from Word document
        try:
            from docx import Document
        except ImportError:
            raise ValueError("DOCX support not available - install python-docx: pip install python-docx") from None

        try:
            doc = Document(BytesIO(file_data))
            text = "\n\n".join([paragraph.text for paragraph in doc.paragraphs if paragraph.text.strip()])
            if not text.strip():
                raise ValueError("Document appears to be empty or contains no readable text")
            return text
        except ValueError:
            raise
        except Exception as e:
            logger.error(f"Error extracting text from .docx file: {str(e)}")
            raise ValueError(f"Failed to extract text from Word document: {str(e)}") from e

    elif file_ext == "xlsx":
        # Extract text from Excel spreadsheet
        if not XLSX_SUPPORT:
            raise ValueError("XLSX support not available - install openpyxl: pip install openpyxl")
        try:
            workbook = load_workbook(BytesIO(file_data), data_only=True)
            text_parts = []
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                text_parts.append(f"=== Sheet: {sheet_name} ===\n")
                for row in sheet.iter_rows(values_only=True):
                    row_text = "\t".join([str(cell) if cell is not None else "" for cell in row])
                    if row_text.strip():
                        text_parts.append(row_text)
            text = "\n".join(text_parts)
            if not text.strip():
                raise ValueError("Spreadsheet appears to be empty")
            return text
        except Exception as e:
            logger.error(f"❌ Error extracting text from .xlsx file: {str(e)}")
            raise ValueError(f"Failed to extract text from Excel file: {str(e)}") from e

    elif file_ext == "pptx":
        # Extract text from PowerPoint presentation
        if not PPTX_SUPPORT:
            raise ValueError("PPTX support not available - install python-pptx: pip install python-pptx")
        try:
            prs = Presentation(BytesIO(file_data))
            text_parts = []
            for slide_num, slide in enumerate(prs.slides, 1):
                text_parts.append(f"=== Slide {slide_num} ===\n")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_parts.append(shape.text)
            text = "\n\n".join(text_parts)
            if not text.strip():
                raise ValueError("Presentation appears to be empty")
            return text
        except Exception as e:
            logger.error(f"❌ Error extracting text from .pptx file: {str(e)}")
            raise ValueError(f"Failed to extract text from PowerPoint file: {str(e)}") from e

    elif file_ext == "pdf":
        # Extract text from PDF - try text extraction first, then OCR fallback
        try:
            from pypdf import PdfReader

            pdf_reader = PdfReader(BytesIO(file_data))
            text_parts = []
            for page_num, page in enumerate(pdf_reader.pages, 1):
                text = page.extract_text()
                if text.strip():
                    text_parts.append(f"=== Page {page_num} ===\n{text}")
            full_text = "\n\n".join(text_parts)

            # If we got meaningful text, return it
            if full_text.strip() and len(full_text.strip()) > 50:
                return full_text

            # Otherwise, try OCR if available
            if OCR_SUPPORT:
                logger.info(f"   PDF has no extractable text, attempting OCR for {filename}...")
                return _extract_pdf_with_ocr(file_data, filename)
            else:
                if not full_text.strip():
                    raise ValueError(
                        "PDF appears to be image-based and OCR is not available. "
                        "Install pytesseract and pdf2image for OCR support."
                    )
                return full_text

        except ImportError:
            raise ValueError("PDF support not available - install pypdf: pip install pypdf") from None
        except Exception as e:
            # If text extraction failed, try OCR as last resort
            if OCR_SUPPORT and "image-based" not in str(e):
                try:
                    logger.info(f"   Text extraction failed, attempting OCR for {filename}...")
                    return _extract_pdf_with_ocr(file_data, filename)
                except Exception as ocr_error:
                    logger.error(f"❌ OCR also failed for {filename}: {str(ocr_error)}")
            logger.error(f"❌ Error extracting text from .pdf file: {str(e)}")
            raise ValueError(f"Failed to extract text from PDF: {str(e)}") from e

    elif file_ext == "txt":
        # Decode text file
        try:
            return file_data.decode("utf-8")
        except UnicodeDecodeError as e:
            raise ValueError(f"Failed to decode text file: {str(e)}") from None

    elif file_ext == "rtf":
        # Extract text from RTF file
        try:
            from striprtf.striprtf import rtf_to_text
        except ImportError:
            raise ValueError("RTF support not available - install striprtf: pip install striprtf") from None

        try:
            # RTF files are text-based, decode first
            rtf_content = file_data.decode("utf-8", errors="replace")
            text = rtf_to_text(rtf_content)
            if not text.strip():
                raise ValueError("RTF document appears to be empty")
            return text
        except ValueError:
            raise
        except Exception as e:
            logger.error(f"❌ Error extracting text from .rtf file: {str(e)}")
            raise ValueError(f"Failed to extract text from RTF document: {str(e)}") from e

    else:
        # Try to decode as text (fallback)
        try:
            return file_data.decode("utf-8")
        except UnicodeDecodeError:
            raise ValueError(
                f"Unsupported file type: .{file_ext}. Supported formats: .pdf, .txt, .docx, .xlsx, .pptx, .rtf"
            ) from None


def chunk_text(
    text: str,
    chunk_size: int = TEXT_CHUNKING.DEFAULT_CHUNK_SIZE,
    overlap: int = TEXT_CHUNKING.DEFAULT_OVERLAP,
) -> List[Dict]:
    """Split text into overlapping chunks.

    Args:
        text: The text to chunk
        chunk_size: Target size of each chunk in characters
        overlap: Number of characters to overlap between chunks

    Returns:
        List of dicts with 'content' and 'chunk_index'
    """
    chunks = []
    start = 0
    chunk_index = 0

    while start < len(text):
        # Get chunk
        end = start + chunk_size
        chunk_text = text[start:end]

        # If not the last chunk, try to break at a sentence or word boundary
        if end < len(text):
            # Look for sentence ending
            last_period = chunk_text.rfind(".")
            last_question = chunk_text.rfind("?")
            last_exclamation = chunk_text.rfind("!")

            boundary = max(last_period, last_question, last_exclamation)

            # If found a sentence boundary in the last X% of chunk, use it
            if boundary > chunk_size * TEXT_CHUNKING.SENTENCE_BOUNDARY_THRESHOLD:
                end = start + boundary + 1
                chunk_text = text[start:end]

        # Only add non-empty chunks
        if chunk_text.strip():
            chunks.append({"content": chunk_text.strip(), "chunk_index": chunk_index})
            chunk_index += 1

        # Move start position (with overlap)
        start = end - overlap

    return chunks


@retry_with_backoff(max_retries=3, initial_delay=2.0, max_delay=30.0, exponential_base=2.0)
def _call_voyage_embed_api(batch: List[str], model: str, input_type: str) -> List[List[float]]:
    """Internal function to call Voyage AI API with retry logic.

    This function is decorated with retry logic to handle transient failures
    such as network issues, rate limits, and temporary API unavailability.

    Args:
        batch: List of texts to embed
        model: Model name to use
        input_type: Type of input (document or query)

    Returns:
        List of embedding vectors

    Raises:
        EmbeddingGenerationError: If all retry attempts fail
    """
    try:
        result = vo.embed(batch, model=model, input_type=input_type)
        return result.embeddings
    except Exception as e:
        # Wrap external service errors with our custom exception
        raise wrap_external_service_error(e, "Voyage AI", "embedding generation") from None


def generate_embeddings(
    texts: List[str],
    batch_size: int = EMBEDDING.DEFAULT_BATCH_SIZE,
    input_type: str = EMBEDDING.INPUT_TYPE_DOCUMENT,
) -> List[List[float]]:
    """Generate embeddings for a list of texts using Voyage AI.

    Handles batching to respect VoyageAI API limits:
    - Max 1000 texts per batch
    - Max 320,000 tokens per batch (~1.28M characters assuming 4 chars/token).

    Args:
        texts: List of text strings to embed
        batch_size: Number of texts to process per batch (default from EMBEDDING config)
        input_type: Type of input - "document" for indexing, "query" for search queries

    Returns:
        List of embedding vectors (1024 dimensions each)
    """
    all_embeddings = []
    total_batches = (len(texts) + batch_size - 1) // batch_size

    logger.info(
        f"   Processing {len(texts)} texts in {total_batches} batches "
        f"(batch_size={batch_size}, input_type={input_type})..."
    )

    try:
        for i in range(0, len(texts), batch_size):
            batch_num = (i // batch_size) + 1
            batch = texts[i : i + batch_size]

            logger.info(f"      Processing batch {batch_num}/{total_batches} ({len(batch)} texts)...")

            # Call Voyage AI API with retry logic
            embeddings = _call_voyage_embed_api(batch=batch, model=EMBEDDING.MODEL_NAME, input_type=input_type)
            all_embeddings.extend(embeddings)

        logger.info(f"   ✓ Successfully generated {len(all_embeddings)} embeddings")
        return all_embeddings

    except Exception as e:
        logger.error(f"❌ Error generating embeddings: {str(e)}")
        raise EmbeddingGenerationError(
            f"Failed to generate embeddings after retries: {str(e)}",
            details={"batch_count": total_batches, "text_count": len(texts)},
        ) from None


def process_document(document_id: str) -> Dict:
    """Process a document: extract text, chunk, embed, store.

    Args:
        document_id: UUID of the document to process

    Returns:
        Processing result summary
    """
    logger.info(f"\nProcessing document: {document_id}")

    try:
        # Get document metadata
        document = documents_repo.get_document(document_id)

        if not document:
            raise ValueError(f"Document {document_id} not found")

        logger.info(f"   Filename: {document['filename']}")

        # Set status to processing
        documents_repo.update_document(document_id, {"processing_status": "processing"})

        # Download file from storage
        # TODO: PocketBase file API migration pending (Plan 3).
        # For now, reconstruct the storage path and use PocketBase file download.
        storage_path = document["storage_url"].split("/documents/")[-1]
        # TODO: Replace with PocketBase file download once file storage is migrated
        raise NotImplementedError(
            f"File download not yet migrated to PocketBase. storage_path={storage_path}"
        )

        # Extract text based on file type
        text = extract_text_from_file(file_data, document["filename"])
        logger.info(f"   Text length: {len(text)} characters")

        # Chunk text
        chunks = chunk_text(text, chunk_size=800, overlap=200)
        logger.info(f"   Created {len(chunks)} chunks")

        # Extract text content for embedding
        chunk_texts = [chunk["content"] for chunk in chunks]

        # Generate embeddings
        logger.info("   Generating embeddings...")
        embeddings = generate_embeddings(chunk_texts)
        logger.info(f"   ✓ Generated {len(embeddings)} embeddings")

        # Store chunks in database (without embeddings) and upsert vectors to Pinecone
        logger.info("   Storing chunks in database...")

        # Collect all chunk data first
        chunks_to_insert = []
        pinecone_vectors = []
        for _i, (chunk, embedding) in enumerate(zip(chunks, embeddings, strict=False)):
            # Sanitize content: remove null bytes which PostgreSQL cannot handle
            sanitized_content = chunk["content"].replace("\x00", "")

            chunk_data = {
                "document_id": document_id,
                "content": sanitized_content,
                "chunk_index": chunk["chunk_index"],
                "metadata": {
                    "filename": document["filename"],
                    "chunk_size": len(sanitized_content),
                },
            }
            chunks_to_insert.append(chunk_data)

        # Batch insert chunks into PocketBase
        BATCH_SIZE = DATABASE.BATCH_INSERT_SIZE
        stored_count = 0
        all_inserted_chunks = []

        for batch_start in range(0, len(chunks_to_insert), BATCH_SIZE):
            batch = chunks_to_insert[batch_start : batch_start + BATCH_SIZE]
            for chunk_data in batch:
                inserted = documents_repo.create_document_chunk(chunk_data)
                all_inserted_chunks.append(inserted)
                stored_count += 1

        # Upsert vectors to Pinecone using the chunk IDs from PocketBase
        for inserted_chunk, embedding in zip(all_inserted_chunks, embeddings, strict=False):
            pinecone_vectors.append(
                {
                    "id": str(inserted_chunk["id"]),
                    "values": embedding,
                    "metadata": {
                        "document_id": document_id,
                        "chunk_index": inserted_chunk.get("chunk_index", 0),
                        "source_type": "document",
                    },
                }
            )

        if pinecone_vectors:
            from services.pinecone_service import upsert_vectors

            upsert_vectors(vectors=pinecone_vectors, namespace="document_chunks")

        logger.info(
            f"   Stored {stored_count} chunks in {(len(chunks_to_insert) + BATCH_SIZE - 1) // BATCH_SIZE} batch(es)"
        )

        # Mark document as processed
        documents_repo.update_document(document_id, {
            "processed": True, "processing_status": "completed", "processing_error": None,
        })

        logger.info("   ✓ Marked document as processed")

        return {
            "document_id": document_id,
            "filename": document["filename"],
            "chunks_created": len(chunks),
            "chunks_stored": stored_count,
            "text_length": len(text),
            "status": "completed",
        }

    except DocumentProcessingError as e:
        # Handle known document processing errors
        error_message = e.message if hasattr(e, "message") else str(e)
        logger.error(
            f"   Document processing error: {error_message}",
            extra={
                "document_id": document_id,
                "error_type": e.__class__.__name__,
                "error_code": e.error_code if hasattr(e, "error_code") else None,
            },
        )

        # Update document with error status
        try:
            documents_repo.update_document(document_id, {
                "processing_status": "failed",
                "processing_error": error_message,
                "processed": True,
            })
            logger.info("   Marked document as failed")
        except Exception as update_error:
            logger.error(f"   Failed to update error status: {update_error}")

        # Don't re-raise - background task should not fail
        return {"document_id": document_id, "status": "failed", "error": error_message}

    except Exception as e:
        # Handle unexpected errors
        error_message = f"Unexpected error: {str(e)}"
        logger.exception("   Unexpected processing error", extra={"document_id": document_id}, exc_info=True)

        # Update document with error status
        try:
            documents_repo.update_document(document_id, {
                "processing_status": "failed",
                "processing_error": error_message,
                "processed": True,
            })
            logger.info("   Marked document as failed")
        except Exception as update_error:
            logger.error(f"   Failed to update error status: {update_error}")

        # Don't re-raise - background task should not fail
        return {"document_id": document_id, "status": "failed", "error": error_message}


async def process_document_with_classification(
    document_id: str, auto_classify: bool = True, anthropic_client=None
) -> Dict:
    """Process a document with optional auto-classification for agent relevance.

    This is an async wrapper around process_document that adds automatic
    agent classification after the document is processed and chunked.

    Args:
        document_id: UUID of the document to process
        auto_classify: Whether to run auto-classification (default: True)
        anthropic_client: Optional Anthropic client for LLM classification

    Returns:
        Processing result summary including classification info
    """
    import asyncio

    # First, run the standard document processing
    result = await asyncio.to_thread(process_document, document_id)

    if result.get("status") != "completed":
        return result

    # If classification is disabled or processing failed, skip
    if not auto_classify:
        return result

    try:
        # Get the first 3 chunks for classification
        esc_doc_id = pb.escape_filter(document_id)
        chunk_list_result = pb.list_records(
            "document_chunks",
            filter=f"document_id='{esc_doc_id}'",
            sort="chunk_index",
            per_page=3,
        )
        chunk_items = chunk_list_result.get("items", [])

        if not chunk_items:
            logger.warning(f"No chunks found for document {document_id}, skipping classification")
            result["classification"] = {"status": "skipped", "reason": "no_chunks"}
            return result

        sample_chunks = [chunk["content"] for chunk in chunk_items]

        # Import classifier here to avoid circular imports
        from services.document_classifier import classify_document_for_agents

        # Run classification
        classification_result = await classify_document_for_agents(
            document_id=document_id,
            chunks=sample_chunks,
            anthropic_client=anthropic_client,
            auto_store=True,
        )

        # Add classification info to result
        result["classification"] = {
            "status": "completed",
            "method": classification_result.classification_method,
            "requires_review": classification_result.requires_user_review,
            "review_reason": classification_result.review_reason,
            "detected_type": classification_result.detected_type,
            "agents": [
                {"name": c.agent_name, "confidence": c.confidence} for c in classification_result.classifications
            ],
            "processing_time_ms": classification_result.processing_time_ms,
        }

        logger.info(
            f"Document {document_id} classified: "
            f"{len(classification_result.classifications)} agents, "
            f"review_needed={classification_result.requires_user_review}"
        )

    except Exception as e:
        logger.error(f"Classification failed for document {document_id}: {e}")
        result["classification"] = {"status": "failed", "error": str(e)}

    # Generate digest for the document
    try:
        from services.document_digests import generate_and_store_digest

        digest = await generate_and_store_digest(document_id)
        result["digest"] = {"status": "completed" if digest else "skipped"}
        if digest:
            logger.info(f"Digest generated for document {document_id} ({len(digest)} chars)")
    except Exception as e:
        logger.warning(f"Digest generation failed for document {document_id}: {e}")
        result["digest"] = {"status": "failed", "error": str(e)}

    # Check if this document affects any opportunities and update justifications
    try:
        from services.project_kb_sync import check_and_sync_opportunities_for_document

        sync_result = await check_and_sync_opportunities_for_document(document_id)
        result["opportunity_sync"] = {
            "status": "completed",
            "opportunities_updated": sync_result.get("opportunities_updated", 0) if sync_result else 0,
        }
        logger.info(f"Opportunity sync completed for document {document_id}")
    except Exception as e:
        logger.warning(f"Opportunity sync failed for document {document_id}: {e}")
        result["opportunity_sync"] = {"status": "failed", "error": str(e)}

    # Auto-extract potential tasks for Taskmaster (only from meeting summaries)
    # Check if document is in the meeting summaries folder
    # Note: Path is case-sensitive, matches Obsidian vault structure
    MEETING_SUMMARIES_FOLDER = "Granola/Meeting-summaries"
    try:
        doc_record = documents_repo.get_document(document_id)

        obsidian_path = doc_record.get("obsidian_file_path", "") if doc_record else ""
        is_meeting_summary = obsidian_path.startswith(MEETING_SUMMARIES_FOLDER)

        if is_meeting_summary:
            from services.task_auto_extractor import extract_tasks_from_document

            task_result = await extract_tasks_from_document(document_id=document_id, auto_store=True)
            result["task_extraction"] = {
                "status": task_result.get("status", "unknown"),
                "tasks_found": task_result.get("tasks_found", 0),
                "tasks_stored": task_result.get("tasks_stored", 0),
            }
            if task_result.get("tasks_found", 0) > 0:
                logger.info(
                    f"Extracted {task_result['tasks_found']} potential tasks from meeting summary {document_id}"
                )
        else:
            result["task_extraction"] = {"status": "skipped", "reason": "not_meeting_summary"}
            logger.debug(f"Skipping task extraction for non-meeting-summary document {document_id}")
    except Exception as e:
        logger.warning(f"Task extraction failed for document {document_id}: {e}")
        result["task_extraction"] = {"status": "failed", "error": str(e)}

    return result


def process_conversation_to_kb(conversation_id: str, user_id: str) -> Dict:
    """Process a conversation into chunks and add to knowledge base.

    Args:
        conversation_id: UUID of the conversation
        user_id: UUID of the user (for access control)

    Returns:
        Processing result summary
    """
    logger.info(f"\nProcessing conversation to KB: {conversation_id}")

    # Get conversation metadata
    conversation = conversations_repo.get_conversation(conversation_id)

    if not conversation:
        raise ValueError(f"Conversation {conversation_id} not found or access denied")

    logger.info(f"   Title: {conversation.get('title', 'Untitled')}")

    # Check if conversation is already in knowledge base
    if conversation.get("in_knowledge_base"):
        logger.warning("   Conversation already in knowledge base")
        raise ValueError(f"Conversation '{conversation.get('title', 'Untitled')}' is already in your knowledge base")

    # Get all messages from the conversation
    messages = conversations_repo.get_conversation_messages(conversation_id)

    if not messages:
        raise ValueError(f"No messages found in conversation {conversation_id}")
    logger.info(f"   Found {len(messages)} messages")

    # Combine messages into conversation text with formatting
    conversation_text = f"Conversation: {conversation.get('title', 'Untitled')}\n\n"
    for msg in messages:
        role_label = "User" if msg["role"] == "user" else "Assistant"
        conversation_text += f"{role_label}: {msg['content']}\n\n"

    logger.info(f"   Total conversation length: {len(conversation_text)} characters")

    # Chunk the conversation
    chunks = chunk_text(conversation_text, chunk_size=800, overlap=200)
    logger.info(f"   Created {len(chunks)} chunks")

    # Extract text content for embedding
    chunk_texts = [chunk["content"] for chunk in chunks]

    # Generate embeddings
    logger.info("   Generating embeddings...")
    embeddings = generate_embeddings(chunk_texts)
    logger.info(f"   ✓ Generated {len(embeddings)} embeddings")

    # Calculate embedding cost (Voyage AI voyage-3: ~$0.02 per 1M tokens)
    # Rough estimate: 4 characters ≈ 1 token for English
    total_chars = sum(len(text) for text in chunk_texts)
    estimated_tokens = total_chars / 4
    embedding_cost_usd = (estimated_tokens / 1_000_000) * 0.02
    logger.info(f"   💰 Embedding cost: ~${embedding_cost_usd:.6f} ({estimated_tokens:.0f} tokens)")

    # Store chunks with embeddings in database with transaction safety
    logger.info("   Storing chunks in database...")

    try:
        # Collect all chunk data first
        chunks_to_insert = []
        for _i, (chunk, embedding) in enumerate(zip(chunks, embeddings, strict=False)):
            # Sanitize content: remove null bytes which PostgreSQL cannot handle
            sanitized_content = chunk["content"].replace("\x00", "")

            chunk_data = {
                "conversation_id": conversation_id,
                "source_type": "conversation",
                "content": sanitized_content,
                "chunk_index": chunk["chunk_index"],
                "metadata": {
                    "conversation_title": conversation.get("title", "Untitled"),
                    "message_count": len(messages),
                    "chunk_size": len(sanitized_content),
                },
            }
            chunks_to_insert.append(chunk_data)

        # Batch insert chunks into PocketBase
        BATCH_SIZE = DATABASE.BATCH_INSERT_SIZE
        stored_count = 0
        all_inserted_chunks = []

        for batch_start in range(0, len(chunks_to_insert), BATCH_SIZE):
            batch = chunks_to_insert[batch_start : batch_start + BATCH_SIZE]
            for chunk_data in batch:
                inserted = documents_repo.create_document_chunk(chunk_data)
                all_inserted_chunks.append(inserted)
                stored_count += 1

        # Upsert vectors to Pinecone
        pinecone_vectors = []
        for inserted_chunk, embedding in zip(all_inserted_chunks, embeddings, strict=False):
            pinecone_vectors.append(
                {
                    "id": str(inserted_chunk["id"]),
                    "values": embedding,
                    "metadata": {
                        "conversation_id": conversation_id,
                        "chunk_index": inserted_chunk.get("chunk_index", 0),
                        "source_type": "conversation",
                    },
                }
            )

        if pinecone_vectors:
            from services.pinecone_service import upsert_vectors

            upsert_vectors(vectors=pinecone_vectors, namespace="document_chunks")

        logger.info(
            f"   Stored {stored_count} chunks in {(len(chunks_to_insert) + BATCH_SIZE - 1) // BATCH_SIZE} batch(es)"
        )

        # Only mark conversation as in knowledge base if ALL chunks succeeded
        from datetime import datetime, timezone
        conversations_repo.update_conversation(conversation_id, {
            "in_knowledge_base": True,
            "added_to_kb_at": datetime.now(timezone.utc).isoformat(),
        })

        logger.info("   ✓ Marked conversation as in knowledge base")

        return {
            "conversation_id": conversation_id,
            "title": conversation.get("title", "Untitled"),
            "chunks_created": len(chunks),
            "chunks_stored": stored_count,
            "message_count": len(messages),
            "embedding_cost_usd": round(embedding_cost_usd, 6),
            "estimated_tokens": int(estimated_tokens),
            "status": "completed",
        }

    except Exception as e:
        # Rollback: Clean up any partially inserted chunks
        logger.error(f"   Error during KB processing: {str(e)}")
        logger.info(f"   Rolling back {stored_count} chunks...")

        try:
            esc_conv_id = pb.escape_filter(conversation_id)
            chunks_to_delete = pb.get_all(
                "document_chunks",
                filter=f"conversation_id='{esc_conv_id}' && source_type='conversation'",
            )
            for chunk in chunks_to_delete:
                pb.delete_record("document_chunks", chunk["id"])
            logger.info("   Rollback complete")
        except Exception as rollback_error:
            logger.error(f"   Rollback failed: {rollback_error}")

        # Re-raise the original error with context
        raise ValueError(f"Failed to add conversation to knowledge base: {str(e)}") from None


def remove_conversation_from_kb(conversation_id: str, user_id: str) -> Dict:
    """Remove a conversation from the knowledge base.

    Args:
        conversation_id: UUID of the conversation
        user_id: UUID of the user (for access control)

    Returns:
        Removal result summary
    """
    logger.info(f"\nRemoving conversation from KB: {conversation_id}")

    # Verify conversation exists
    conversation = conversations_repo.get_conversation(conversation_id)

    if not conversation:
        raise ValueError(f"Conversation {conversation_id} not found or access denied")

    # Get chunk IDs before deleting (for Pinecone cleanup)
    esc_conv_id = pb.escape_filter(conversation_id)
    chunks_to_delete = pb.get_all(
        "document_chunks",
        filter=f"conversation_id='{esc_conv_id}'",
    )
    chunk_ids_to_delete = [str(c["id"]) for c in chunks_to_delete]

    # Delete from Pinecone
    if chunk_ids_to_delete:
        from services.pinecone_service import delete_vectors

        delete_vectors(ids=chunk_ids_to_delete, namespace="document_chunks")

    # Delete all chunks for this conversation
    chunks_deleted = 0
    for chunk in chunks_to_delete:
        pb.delete_record("document_chunks", chunk["id"])
        chunks_deleted += 1

    logger.info(f"   Deleted {chunks_deleted} chunks")

    # Mark conversation as not in knowledge base
    conversations_repo.update_conversation(conversation_id, {
        "in_knowledge_base": False,
        "added_to_kb_at": None,
    })

    logger.info("   Marked conversation as removed from knowledge base")

    return {
        "conversation_id": conversation_id,
        "chunks_deleted": chunks_deleted,
        "status": "removed",
    }


def preprocess_query(query: str) -> str:
    """Preprocess search query to improve relevance by removing stop words.

    and focusing on key terms.

    Args:
        query: The raw search query

    Returns:
        Preprocessed query with stop words removed
    """
    # Common stop words that dilute semantic meaning
    stop_words = {
        "what",
        "are",
        "is",
        "the",
        "a",
        "an",
        "and",
        "or",
        "but",
        "in",
        "on",
        "at",
        "to",
        "for",
        "of",
        "with",
        "by",
        "from",
        "as",
        "this",
        "that",
        "these",
        "those",
        "i",
        "you",
        "we",
        "they",
        "it",
        "my",
        "your",
        "our",
        "their",
        "can",
        "could",
        "would",
        "should",
        "will",
        "do",
        "does",
        "did",
        "have",
        "has",
        "had",
        "be",
        "been",
        "am",
        "was",
        "were",
        "me",
        "us",
        "them",
        "there",
        "here",
        "where",
        "when",
        "how",
        "why",
        "which",
        "who",
        "whom",
        "whose",
        "tell",
        "show",
        "give",
        "find",
        "get",
        # Meta-terms that don't add semantic value to queries
        "any",
        "some",
        "about",
        "regarding",
        "information",
        "info",
        "data",
        "details",
        "know",
        "help",
        "need",
        "want",
        "please",
    }

    # Split query into words and filter out stop words
    words = query.lower().split()
    key_terms = [word.strip(".,?!:;") for word in words if word.lower() not in stop_words]

    # If all words were stop words, return original query to avoid empty search
    if not key_terms:
        return query

    # Reconstruct query from key terms
    preprocessed = " ".join(key_terms)

    logger.info("   Query preprocessing:")
    logger.info(f"      Original: '{query}'")
    logger.info(f"      Preprocessed: '{preprocessed}'")

    return preprocessed


def detect_query_type(query: str) -> str:
    """Detect the type of query to determine appropriate search strategy.

    Args:
        query: The search query

    Returns:
        Query type:
        - 'document_reference' for queries explicitly referencing uploaded documents
        - 'factual' for direct questions
        - 'exploratory' for general queries
    """
    query_lower = query.lower().strip()

    # HIGHEST PRIORITY: Detect document-reference queries
    # These should use very permissive thresholds since user is explicitly asking about their uploads
    document_reference_patterns = [
        "this document",
        "that document",
        "the document",
        "this file",
        "that file",
        "the file",
        "this pdf",
        "that pdf",
        "the pdf",
        "this report",
        "that report",
        "the report",
        "i uploaded",
        "i just uploaded",
        "document i uploaded",
        "file i uploaded",
        "uploaded document",
        "uploaded file",
        "my document",
        "my file",
        "my upload",
        "tell me about this",
        "analyze this",
        "summarize this",
        "what does this say",
        "what is in this",
        # Meeting/transcript references - user is asking about their meeting content
        "my meeting",
        "my meetings",
        "the meeting",
        "recent meeting",
        "recent meetings",
        "my transcript",
        "my transcripts",
        "the transcript",
        "transcript",
        "action item",
        "action items",
        "action points",
        "meeting notes",
        "meeting note",
        "from the meeting",
        "from my meeting",
        "what was discussed",
        "what did we discuss",
        "what we discussed",
        "follow up",
        "follow-up",
        "followup",
        "next steps",
        "decisions made",
        "key decisions",
        "takeaways",
        "key takeaways",
    ]

    for pattern in document_reference_patterns:
        if pattern in query_lower:
            return "document_reference"

    # Direct question words indicate a need for specific factual answers
    question_words = ["what", "who", "where", "when", "how", "why", "which", "whose"]

    # Check if query starts with a question word
    for qword in question_words:
        if query_lower.startswith(qword + " "):
            return "factual"

    # Questions with "is/are/do/does" patterns
    question_patterns = [
        "is there",
        "are there",
        "is the",
        "are the",
        "do we",
        "does the",
        "can you",
        "could you",
        "will the",
        "would the",
        "do you",
        "did you",
        "have you",
        "has there",
        "does it",
    ]

    for pattern in question_patterns:
        if query_lower.startswith(pattern):
            return "factual"

    # Default to exploratory for general statements or broad queries
    return "exploratory"


def search_similar_chunks(
    query: str,
    client_id: str,
    limit: int = 5,
    include_conversations: bool = True,
    min_similarity: float = 0.0,
    user_id: str = None,
    document_ids: List[str] = None,
    conversation_id: str = None,
    agent_ids: List[str] = None,
    fallback_threshold: int = 3,
) -> List[Dict]:
    """Search for document chunks similar to a query.

    Uses Redis cache for identical queries (1-hour TTL) to avoid expensive
    vector similarity searches.

    Args:
        query: The search query
        client_id: Client ID to filter results
        limit: Maximum number of results
        include_conversations: Whether to include conversation chunks (default: True)
        min_similarity: Minimum similarity score to include results (default: 0.0)
        user_id: Optional user ID to filter results
        document_ids: Optional list of document IDs to search within (if None, searches all documents)
        conversation_id: Optional conversation ID to prioritize files referenced in conversation
        agent_ids: Optional list of agent IDs to prioritize documents tagged for these agents
        fallback_threshold: Minimum chunks needed before falling back to all docs (default: 3)

    Returns:
        List of similar chunks with similarity scores
    """
    logger.info(f"\n🔍 Searching for: '{query}'")
    logger.info(f"   Client: {client_id}")
    logger.info(f"   Include conversations: {include_conversations}")
    logger.info(f"   Min similarity: {min_similarity}")
    logger.info(f"   Conversation context: {conversation_id}")
    if document_ids:
        logger.info(f"   Filtering by document IDs: {document_ids}")
    if agent_ids:
        logger.info(f"   Filtering by agent IDs: {agent_ids} (fallback threshold: {fallback_threshold})")

    # Try to get from cache first (1-hour TTL)
    # Cache key includes all parameters to ensure exact match
    agent_ids_key = ":".join(sorted(agent_ids)) if agent_ids else "none"
    doc_ids_key = ":".join(sorted(document_ids)) if document_ids else "none"
    cache_key_suffix = (
        f"{query}:{limit}:{include_conversations}:{min_similarity}:{conversation_id}:{agent_ids_key}:{doc_ids_key}"
    )
    cached_results = get_cached_search_results(cache_key_suffix, client_id)
    if cached_results is not None:
        logger.info(f"   ✅ Search results loaded from cache ({len(cached_results)} results)")
        return cached_results

    logger.debug("   📋 Cache miss - performing vector search")

    # Get conversation-referenced document IDs if conversation_id provided
    conversation_doc_ids = []
    if conversation_id:
        try:
            # Query all messages in the conversation
            conv_messages = conversations_repo.get_conversation_messages(conversation_id)

            if conv_messages:
                message_ids = [msg["id"] for msg in conv_messages]

                # Get all documents referenced in this conversation
                all_doc_ids = set()
                for mid in message_ids:
                    msg_docs = conversations_repo.get_message_documents(mid)
                    for md in msg_docs:
                        all_doc_ids.add(md["document_id"])

                if all_doc_ids:
                    conversation_doc_ids = list(all_doc_ids)
                    logger.info(f"   Found {len(conversation_doc_ids)} documents referenced in conversation")
        except Exception as e:
            logger.warning(f"   Failed to fetch conversation documents: {e}")

    # Detect query type for adaptive strategy
    query_type = detect_query_type(query)
    logger.info(f"   Query type: {query_type}")

    # Preprocess query to extract key terms
    preprocessed_query = preprocess_query(query)

    # Adjust minimum similarity based on query type
    # Factual questions need higher precision, exploratory queries can be more permissive
    # Document reference queries are very permissive since user explicitly mentioned their uploads
    if min_similarity == 0.0:  # Only adjust if not explicitly set
        if query_type == "document_reference":
            min_similarity = SEARCH.SIMILARITY_THRESHOLDS["document_reference"]
            logger.info(f"   Adjusted min_similarity to {min_similarity} for document reference query (permissive)")
        elif query_type == "factual":
            min_similarity = SEARCH.SIMILARITY_THRESHOLDS["factual"]
            logger.info(f"   Adjusted min_similarity to {min_similarity} for factual query")
        else:
            min_similarity = SEARCH.SIMILARITY_THRESHOLDS["exploratory"]
            logger.info(f"   Adjusted min_similarity to {min_similarity} for exploratory query")

    # Generate embedding for query using "query" input type for better search relevance
    query_embedding = generate_embeddings([preprocessed_query], input_type=EMBEDDING.INPUT_TYPE_QUERY)[0]

    # Detect if query is about meetings/transcripts/action items - use document_type filtering
    query_lower = query.lower()
    meeting_action_keywords = [
        "meeting",
        "transcript",
        "action item",
        "action items",
        "follow up",
        "follow-up",
        "next steps",
        "takeaways",
        "decisions",
        "discussed",
        "talked about",
        "what did we",
        "what was discussed",
        "from the meeting",
        "from my meeting",
    ]
    # Also detect "what should I work on" style queries - these want recent action items
    work_task_keywords = [
        "what should i work on",
        "what do i need to do",
        "my tasks",
        "my to do",
        "my priorities",
        "what's on my plate",
        "what am i working on",
        "today's tasks",
        "this week",
        "my deliverables",
        "my assignments",
    ]
    # Detect recency keywords - user wants RECENT content only
    recency_keywords = [
        "recent",
        "last few",
        "latest",
        "this week",
        "past week",
        "last week",
        "today",
        "yesterday",
        "past few days",
        "last couple",
        "most recent",
    ]
    is_meeting_query = any(kw in query_lower for kw in meeting_action_keywords)
    is_work_query = any(kw in query_lower for kw in work_task_keywords)
    is_recency_query = any(kw in query_lower for kw in recency_keywords)

    # Work queries should also search transcripts/notes for action items
    if is_work_query and not is_meeting_query:
        is_meeting_query = True
        logger.info("   Detected work/task query - will search transcripts for action items")

    # Calculate date filter for recency queries (default: last 7 days)
    from datetime import datetime, timedelta, timezone

    recency_date = None
    if is_recency_query or is_work_query:
        recency_date = (datetime.now(timezone.utc) - timedelta(days=7)).isoformat()
        logger.info(f"   Detected recency query - filtering to docs after {recency_date[:10]}")

    # Call vector search function - get more results to prioritize conversation docs
    # Use agent-filtered RPC if agent_ids provided, otherwise use standard RPC
    used_fallback = False
    used_document_type_filter = False

    # Check if Pinecone is available
    from services.pinecone_service import get_index as _get_pinecone_index
    from services.pinecone_service import query_vectors as _query_pinecone

    _use_pinecone = _get_pinecone_index() is not None

    def _pinecone_search_and_hydrate(pc_filter, top_k, threshold):
        """Query Pinecone and hydrate results from PocketBase."""
        matches = _query_pinecone(
            embedding=query_embedding,
            namespace="document_chunks",
            top_k=top_k,
            filter=pc_filter if pc_filter else None,
        )
        # Filter by threshold
        matches = [m for m in matches if m["score"] >= threshold]
        if not matches:
            return []

        # Fetch chunk data from PocketBase by IDs
        chunk_ids = [m["id"] for m in matches]
        scores_map = {m["id"]: m["score"] for m in matches}

        # Build OR filter for batch fetch
        hydrated = []
        for batch_start in range(0, len(chunk_ids), 100):
            batch_ids = chunk_ids[batch_start : batch_start + 100]
            try:
                or_parts = [f"id='{pb.escape_filter(cid)}'" for cid in batch_ids]
                or_filter = " || ".join(or_parts)
                rows = pb.get_all("document_chunks", filter=or_filter)
                for row in rows:
                    row["similarity"] = scores_map.get(str(row["id"]), 0)
                    hydrated.append(row)
            except Exception as e:
                logger.error(f"   Failed to hydrate Pinecone results from PocketBase: {e}")

        # Sort by similarity descending
        hydrated.sort(key=lambda x: x.get("similarity", 0), reverse=True)
        return hydrated

    # First, try document_type filtered search for meeting queries
    if is_meeting_query and not agent_ids:
        logger.info("   Detected meeting/action item query - trying document_type filtered search")
        try:
            if _use_pinecone:
                # Pinecone filter for document types
                pc_filter = {
                    "$and": [
                        {"client_id": {"$eq": client_id}},
                        {"source_type": {"$in": ["transcript", "notes"]}},
                    ]
                }
                if user_id:
                    pc_filter["$and"].append({"user_id": {"$eq": user_id}})

                chunks = _pinecone_search_and_hydrate(pc_filter, limit * 3, min_similarity)

                # Apply recency filter in post-processing if needed
                if recency_date and chunks:
                    chunks = [c for c in chunks if (c.get("created_at") or "") >= recency_date]
            else:
                # Fallback to vec_client search + PocketBase filter
                import vec_client as _vec
                vec_results = _vec.search("document_chunks", preprocessed_query, limit=limit * 3)
                # Filter by source_type and similarity
                chunks = []
                for vr in vec_results:
                    rec_id = vr.get("record_id") or vr.get("id", "")
                    chunk_rec = pb.get_record("document_chunks", rec_id) if rec_id else None
                    if chunk_rec and chunk_rec.get("source_type") in ("transcript", "notes"):
                        sim = vr.get("score", vr.get("similarity", 0))
                        if sim >= min_similarity:
                            chunk_rec["similarity"] = sim
                            chunks.append(chunk_rec)

            # Keep date filter active - don't silently drop it when few results found
            if recency_date and len(chunks) < 2:
                logger.info(
                    f"   Date filter active: only {len(chunks)} recent chunks found (respecting user's recency request)"
                )

            if len(chunks) >= 1:
                used_document_type_filter = True
                logger.info(f"   Found {len(chunks)} chunks from transcripts/notes")

                # For work/task queries, boost recent documents (last 14 days)
                if is_work_query:
                    cutoff = datetime.now(timezone.utc) - timedelta(days=14)
                    for chunk in chunks:
                        chunk_date = chunk.get("created_at")
                        if chunk_date:
                            try:
                                if isinstance(chunk_date, str):
                                    chunk_dt = datetime.fromisoformat(chunk_date.replace("Z", "+00:00"))
                                else:
                                    chunk_dt = chunk_date
                                if chunk_dt > cutoff:
                                    original_sim = chunk.get("similarity", 0)
                                    chunk["similarity"] = min(1.0, original_sim * 1.2)
                                    chunk["recency_boosted"] = True
                            except (ValueError, TypeError):
                                pass
                    chunks.sort(key=lambda x: x.get("similarity", 0), reverse=True)
                    logger.info("   Applied recency boost for work query")
            else:
                logger.info(f"   Only found {len(chunks)} chunks from transcripts/notes, will try broader search")
                chunks = []
        except Exception as rpc_error:
            logger.warning(f"   Document type filtered search failed: {rpc_error}, falling back to standard search")
            chunks = []

    # Only run standard search if we didn't already get results from document_type filter
    if not used_document_type_filter:
        if document_ids:
            logger.info(
                f"   Searching with {len(document_ids)} document IDs, threshold={min_similarity}"
            )
            logger.info(f"   Query embedding length: {len(query_embedding)}")

            try:
                if _use_pinecone:
                    pc_filter = {"document_id": {"$in": document_ids}}
                    chunks = _pinecone_search_and_hydrate(pc_filter, limit * 3, min_similarity)
                else:
                    # Fallback to vec_client search + PocketBase filter by document IDs
                    import vec_client as _vec
                    vec_results = _vec.search("document_chunks", preprocessed_query, limit=limit * 3)
                    chunks = []
                    for vr in vec_results:
                        rec_id = vr.get("record_id") or vr.get("id", "")
                        chunk_rec = pb.get_record("document_chunks", rec_id) if rec_id else None
                        if chunk_rec and chunk_rec.get("document_id") in document_ids:
                            sim = vr.get("score", vr.get("similarity", 0))
                            if sim >= min_similarity:
                                chunk_rec["similarity"] = sim
                                chunks.append(chunk_rec)

                logger.info(f"   Found {len(chunks)} results from specified documents")
                if chunks:
                    logger.info(f"   Top result similarity: {chunks[0].get('similarity', 'N/A')}")
            except Exception as rpc_error:
                logger.error(f"   Document-ID-filtered search failed: {rpc_error}")
                chunks = []
        elif agent_ids:
            logger.info(
                f"   Searching with agent filter, threshold={min_similarity}, agents={agent_ids}"
            )
            logger.info(f"   Query embedding length: {len(query_embedding)}")

            try:
                if _use_pinecone:
                    # For agent-filtered search, first get document IDs tagged for these agents
                    # then search Pinecone with those document IDs
                    try:
                        or_parts = [f"agent_id='{pb.escape_filter(aid)}'" for aid in agent_ids]
                        or_filter = " || ".join(or_parts)
                        agent_doc_records = pb.get_all("document_agents", filter=or_filter)
                        agent_doc_ids = list({d["document_id"] for d in agent_doc_records})
                    except Exception:
                        agent_doc_ids = []

                    if agent_doc_ids:
                        pc_filter = {
                            "$and": [
                                {"client_id": {"$eq": client_id}},
                                {"document_id": {"$in": agent_doc_ids}},
                            ]
                        }
                        if user_id:
                            pc_filter["$and"].append({"user_id": {"$eq": user_id}})
                        chunks = _pinecone_search_and_hydrate(pc_filter, limit * 3, min_similarity)
                    else:
                        chunks = []

                    # Fallback to all docs if agent filter has insufficient results
                    if len(chunks) < fallback_threshold:
                        used_fallback = True
                        logger.info("   Agent filter had insufficient results, using fallback to all documents")
                        pc_filter = {"client_id": {"$eq": client_id}}
                        if user_id:
                            pc_filter["user_id"] = {"$eq": user_id}
                        chunks = _pinecone_search_and_hydrate(pc_filter, limit * 3, min_similarity)
                        for chunk in chunks:
                            chunk["used_fallback"] = True
                else:
                    # Fallback to vec_client search + PocketBase agent filter
                    import vec_client as _vec
                    try:
                        or_parts = [f"agent_id='{pb.escape_filter(aid)}'" for aid in agent_ids]
                        or_filter = " || ".join(or_parts)
                        agent_doc_records = pb.get_all("document_agents", filter=or_filter)
                        agent_doc_ids_set = {d["document_id"] for d in agent_doc_records}
                    except Exception:
                        agent_doc_ids_set = set()

                    vec_results = _vec.search("document_chunks", preprocessed_query, limit=limit * 3)
                    chunks = []
                    for vr in vec_results:
                        rec_id = vr.get("record_id") or vr.get("id", "")
                        chunk_rec = pb.get_record("document_chunks", rec_id) if rec_id else None
                        if chunk_rec:
                            sim = vr.get("score", vr.get("similarity", 0))
                            if sim >= min_similarity:
                                chunk_rec["similarity"] = sim
                                if chunk_rec.get("document_id") in agent_doc_ids_set:
                                    chunks.append(chunk_rec)

                    if len(chunks) < fallback_threshold:
                        used_fallback = True
                        logger.info("   Agent filter had insufficient results, used fallback to all documents")
                        chunks = []
                        for vr in vec_results:
                            rec_id = vr.get("record_id") or vr.get("id", "")
                            chunk_rec = pb.get_record("document_chunks", rec_id) if rec_id else None
                            if chunk_rec:
                                sim = vr.get("score", vr.get("similarity", 0))
                                if sim >= min_similarity:
                                    chunk_rec["similarity"] = sim
                                    chunk_rec["used_fallback"] = True
                                    chunks.append(chunk_rec)

                logger.info(f"   Found {len(chunks)} results (fallback={used_fallback})")
                if chunks:
                    logger.info(f"   Top result similarity: {chunks[0].get('similarity', 'N/A')}")
            except Exception as rpc_error:
                logger.error(f"   Agent-filtered search failed: {rpc_error}")
                chunks = []
        else:
            logger.info(
                f"   Searching with threshold={min_similarity}, "
                f"client_id={client_id}, user_id={user_id}"
            )
            logger.info(f"   Query embedding length: {len(query_embedding)}")

            try:
                if _use_pinecone:
                    pc_filter = {"client_id": {"$eq": client_id}}
                    if user_id:
                        pc_filter["user_id"] = {"$eq": user_id}
                    chunks = _pinecone_search_and_hydrate(pc_filter, limit * 3, min_similarity)
                else:
                    # Fallback to vec_client search
                    import vec_client as _vec
                    vec_results = _vec.search("document_chunks", preprocessed_query, limit=limit * 3)
                    chunks = []
                    for vr in vec_results:
                        rec_id = vr.get("record_id") or vr.get("id", "")
                        chunk_rec = pb.get_record("document_chunks", rec_id) if rec_id else None
                        if chunk_rec:
                            sim = vr.get("score", vr.get("similarity", 0))
                            if sim >= min_similarity:
                                chunk_rec["similarity"] = sim
                                chunks.append(chunk_rec)

                logger.info(f"   Found {len(chunks)} results")
                if chunks:
                    logger.info(f"   Top result similarity: {chunks[0].get('similarity', 'N/A')}")
            except Exception as rpc_error:
                logger.error(f"   Search failed: {rpc_error}")
                chunks = []

    # Filter out conversation chunks if requested
    if not include_conversations:
        chunks = [c for c in chunks if c.get("source_type") != "conversation"]
        logger.info(f"   After filtering: {len(chunks)} document-only results")

    # Filter by document IDs if provided
    if document_ids:
        chunks = [c for c in chunks if c.get("document_id") in document_ids]
        logger.info(f"   After document ID filtering: {len(chunks)} results from specified documents")

    # Filter by minimum similarity score
    if min_similarity > 0.0:
        chunks = [c for c in chunks if c.get("similarity", 0.0) >= min_similarity]
        logger.info(f"   After similarity filtering (>={min_similarity}): {len(chunks)} results")

    # Boost documents matching query keywords (meeting/transcript queries should prioritize those doc types)
    query_lower = query.lower()
    meeting_keywords = ["meeting", "transcript", "discussion", "conversation", "call", "session"]
    if any(kw in query_lower for kw in meeting_keywords):
        for chunk in chunks:
            metadata = chunk.get("metadata", {})
            filename = metadata.get("filename", "").lower()
            title = metadata.get("title", "").lower() if metadata.get("title") else ""

            # Boost chunks from documents with meeting/transcript in filename or title
            if any(kw in filename or kw in title for kw in ["transcript", "meeting", "session", "call"]):
                chunk["similarity"] = min(1.0, chunk.get("similarity", 0.0) * 1.5)  # 50% boost
                chunk["title_boosted"] = True
                logger.debug(f"   Boosted transcript/meeting doc: {filename}")

        # Re-sort after boosting
        chunks.sort(key=lambda x: x.get("similarity", 0.0), reverse=True)
        boosted_count = sum(1 for c in chunks if c.get("title_boosted"))
        if boosted_count > 0:
            logger.info(f"   Boosted {boosted_count} chunks from meeting/transcript documents")

    # Prioritize chunks from conversation-referenced documents
    if conversation_doc_ids:
        # Split into conversation docs and other docs
        conversation_chunks = [c for c in chunks if c.get("document_id") in conversation_doc_ids]
        other_chunks = [c for c in chunks if c.get("document_id") not in conversation_doc_ids]

        # Boost similarity scores for conversation docs (multiply by 1.2)
        for chunk in conversation_chunks:
            chunk["similarity"] = min(1.0, chunk.get("similarity", 0.0) * 1.2)
            chunk["from_conversation"] = True

        # Mark other chunks
        for chunk in other_chunks:
            chunk["from_conversation"] = False

        # Combine with conversation docs first, then re-sort by boosted similarity
        chunks = conversation_chunks + other_chunks
        chunks.sort(key=lambda x: x.get("similarity", 0.0), reverse=True)

        logger.info(f"   Prioritized {len(conversation_chunks)} chunks from conversation documents")

    # Limit to requested number
    chunks = chunks[:limit]

    # Cache results for 1 hour (3600 seconds)
    cache_search_results(cache_key_suffix, client_id, chunks, ttl=3600)

    return chunks
